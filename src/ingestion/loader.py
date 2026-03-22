import logging
import os
import json
import hashlib
from pathlib import Path
from typing import Dict, List, Optional, Any
logger = logging.getLogger(__name__)

class CacheManager:

    def __init__(self, cache_dir: str='data/processed/cache'):
        self.cache_dir = Path(cache_dir)
        self.cache_dir.mkdir(parents=True, exist_ok=True)

    def get_hash(self, file_path: str) -> str:
        p = Path(file_path).resolve()
        stat = p.stat()
        seed = f'{str(p)}|{stat.st_size}|{stat.st_mtime}'
        return hashlib.md5(seed.encode('utf-8')).hexdigest()

    def get_cached(self, file_path: str) -> Optional[Dict]:
        p = self.cache_dir / f'doc_{self.get_hash(file_path)}.json'
        if p.exists():
            try:
                with open(p, 'r', encoding='utf-8') as f:
                    return json.load(f)
            except Exception as e:
                logger.warning(f'Failed to read cache {p}: {e}')
        return None

    def save_cache(self, file_path: str, data: Dict):
        p = self.cache_dir / f'doc_{self.get_hash(file_path)}.json'
        try:
            with open(p, 'w', encoding='utf-8') as f:
                json.dump(data, f, ensure_ascii=False)
        except Exception as e:
            logger.warning(f'Failed to save cache {p}: {e}')

class RawDocument:

    def __init__(self, doc_key: str, source: str, text: str='', image_path: Optional[str]=None, metadata: Optional[Dict]=None):
        self.doc_key = doc_key
        self.source = source
        self.text = text
        self.image_path = image_path
        self.metadata = metadata or {}

    def __repr__(self) -> str:
        return f"RawDocument(key='{self.doc_key}', source='{self.source}', text_len={len(self.text)}, has_image={self.image_path is not None})"

    def to_dict(self) -> Dict:
        return {'doc_key': self.doc_key, 'source': self.source, 'text_length': len(self.text), 'image_path': self.image_path, 'metadata': self.metadata}

def _load_arrow_dataset(dataset_path: str):
    try:
        from datasets import load_from_disk
    except ImportError:
        raise ImportError("The 'datasets' library is required to load Arrow datasets.  Install it with:  pip install datasets")
    return load_from_disk(dataset_path)

def _extract_image(image_obj, save_path: str) -> str:
    os.makedirs(os.path.dirname(save_path), exist_ok=True)
    image_obj.save(save_path, format='PNG')
    return os.path.abspath(save_path)

def _extract_text_from_pdf(file_path: str) -> str:
    try:
        import pdfplumber
    except ImportError:
        raise ImportError("The 'pdfplumber' library is required to extract text from PDFs.  Install it with:  pip install pdfplumber")
    pages_text: List[str] = []
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                pages_text.append(text)
    return '\n\n'.join(pages_text)

def _extract_text_from_docx(file_path: str) -> str:
    try:
        import docx
    except ImportError:
        raise ImportError("The 'python-docx' library is required to extract text from DOCX files.  Install it with:  pip install python-docx")
    try:
        doc = docx.Document(file_path)
    except Exception as exc:
        logger.warning('Failed to open DOCX file %s: %s', file_path, exc)
        return ''
    parts: List[str] = []
    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style_name = (paragraph.style.name or '').lower()
        if 'heading' in style_name:
            parts.append(f'Heading: {text}')
        else:
            parts.append(f'Paragraph:\n{text}')
    for table in doc.tables:
        rows_text: List[str] = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows_text.append(' | '.join(cells))
        if rows_text:
            parts.append('Table:\n' + '\n'.join(rows_text))
    return '\n\n'.join(parts)

def _extract_docx_metadata(file_path: str) -> dict:
    try:
        import docx
        doc = docx.Document(file_path)
        return {'paragraph_count': len([p for p in doc.paragraphs if p.text.strip()]), 'table_count': len(doc.tables)}
    except Exception:
        return {}

def _extract_text_from_pptx(file_path: str) -> str:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        raise ImportError("The 'python-pptx' library is required to extract text from PPTX files.  Install it with:  pip install python-pptx")
    try:
        prs = Presentation(file_path)
    except Exception as exc:
        logger.warning('Failed to open PPTX file %s: %s', file_path, exc)
        return ''
    slides_text: List[str] = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_parts: List[str] = [f'Slide {slide_num}']
        title = ''
        if slide.shapes.title and slide.shapes.title.text.strip():
            title = slide.shapes.title.text.strip()
            slide_parts.append(f'Title: {title}')
        content_lines: List[str] = []
        for shape in slide.shapes:
            if slide.shapes.title and shape.shape_id == slide.shapes.title.shape_id:
                continue
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        content_lines.append(text)
            if shape.has_table:
                rows_text: List[str] = []
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows_text.append(' | '.join(cells))
                if rows_text:
                    content_lines.append('Table:\n' + '\n'.join(rows_text))
        if content_lines:
            slide_parts.append('Content: ' + '\n'.join(content_lines))
        if len(slide_parts) > 1:
            slides_text.append('\n'.join(slide_parts))
    return '\n\n'.join(slides_text)

def _extract_pptx_metadata(file_path: str) -> dict:
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        return {'slide_count': len(prs.slides)}
    except Exception:
        return {}
_FILE_EXTRACTORS = {'.txt': lambda path: Path(path).read_text(encoding='utf-8', errors='replace'), '.pdf': _extract_text_from_pdf, '.docx': _extract_text_from_docx, '.pptx': _extract_text_from_pptx}
_IMAGE_EXTENSIONS = {'.png', '.jpg', '.jpeg', '.tiff', '.tif', '.bmp'}
_VIDEO_EXTENSIONS = {'.mp4', '.avi', '.mkv', '.mov', '.webm', '.flv', '.wmv'}

def _process_video_file(video_path: str, settings: dict=None) -> Optional['RawDocument']:
    settings = settings or {}
    video_settings = settings.get('video', {})
    vp = Path(video_path)
    if not vp.exists():
        logger.error('Video file not found: %s', video_path)
        return None
    output_base = Path(video_settings.get('output_dir', 'data/processed/videos'))
    frames_dir = output_base / 'frames'
    audio_dir = output_base / 'audio'
    video_id = vp.stem
    frame_interval = video_settings.get('frame_interval', 5)
    enable_whisper = video_settings.get('enable_whisper', True)
    enable_ocr = video_settings.get('enable_frame_ocr', True)
    enable_captioning = video_settings.get('enable_frame_captioning', True)
    caption_interval = video_settings.get('caption_interval', 5)
    caption_model = video_settings.get('caption_model', 'Salesforce/blip-image-captioning-base')
    ocr_min_words = video_settings.get('ocr_min_words', 3)
    whisper_model = video_settings.get('whisper_model_size', 'small')
    whisper_device = video_settings.get('whisper_device', 'cpu')
    chunk_interval = video_settings.get('chunk_interval', 30)
    transcript_segments = []
    ocr_results = []
    frame_captions = []
    duration = 0.0
    try:
        from src.video.frame_sampler import FrameSampler
        sampler = FrameSampler(output_dir=str(frames_dir), interval_seconds=frame_interval)
        frames = sampler.sample(str(vp), video_id=video_id)
        duration = sampler.get_video_duration(str(vp))
        logger.info('Sampled %d frames from %s (%.1fs)', len(frames), vp.name, duration)
    except Exception as exc:
        logger.error('Frame sampling failed for %s: %s', vp.name, exc)
        return None
    if enable_ocr and frames:
        try:
            from src.video.frame_ocr import FrameOCR
            frame_ocr = FrameOCR(min_word_count=ocr_min_words)
            ocr_results = frame_ocr.extract_batch(frames)
            logger.info('OCR: %d/%d frames had text', len(ocr_results), len(frames))
        except Exception as exc:
            logger.warning('Frame OCR failed for %s: %s', vp.name, exc)
    if enable_captioning and frames:
        try:
            from src.video.frame_captioner import FrameCaptioner
            captioner = FrameCaptioner(model_name=caption_model, use_openvino=video_settings.get('caption_use_openvino', False))
            frame_captions = captioner.caption_batch(frames, interval=caption_interval)
            logger.info('Captioned %d frames from %s', len(frame_captions), vp.name)
        except ImportError as exc:
            logger.warning('Frame captioning skipped (missing dependency: %s)', exc)
        except Exception as exc:
            logger.warning('Frame captioning failed for %s: %s', vp.name, exc)
    if enable_whisper:
        try:
            from src.video.audio_extractor import AudioExtractor
            from src.video.transcription import WhisperTranscriber
            extractor = AudioExtractor(output_dir=str(audio_dir))
            wav_path = extractor.extract(str(vp))
            if wav_path:
                transcriber = WhisperTranscriber(model_size=whisper_model, device=whisper_device)
                transcript_segments = transcriber.transcribe(wav_path)
                logger.info('Transcribed %d segments from %s', len(transcript_segments), vp.name)
            else:
                logger.info('No audio track in %s — using OCR only', vp.name)
        except ImportError as exc:
            logger.warning('Whisper/moviepy not available (%s) — using OCR only', exc)
        except Exception as exc:
            logger.warning('Transcription failed for %s: %s — using OCR only', vp.name, exc)
    caption_text_parts = []
    for fc in frame_captions:
        caption_text_parts.append(f'Frame at {fc.timestamp:.1f}s: {fc.caption}')
    caption_text = '\n'.join(caption_text_parts)
    frame_paths = [f.frame_path for f in frames] if frames else []
    has_any_text = transcript_segments or ocr_results or caption_text
    if not has_any_text:
        if frame_paths:
            logger.info('No text from OCR/Whisper/captioning for %s — falling back to CLIP visual search on %d frames', vp.name, len(frame_paths))
            return RawDocument(doc_key=video_id, source=str(vp), text=f'[Visual video: {vp.name}]', metadata={'file_type': 'video', 'file_name': vp.name, 'source_directory': str(vp.parent), 'modality': 'video', 'duration': round(duration, 2), 'frame_count': len(frame_paths), 'ocr_frames_with_text': 0, 'transcript_segments': 0, 'caption_count': 0, 'video_frame_paths': frame_paths, 'clip_only': True})
        logger.warning('No text and no frames extracted from %s — skipping', vp.name)
        return None
    try:
        from src.video.video_document_builder import VideoDocumentBuilder
        builder = VideoDocumentBuilder(output_dir=str(output_base), chunk_interval=chunk_interval)
        all_captions = []
        if caption_text:
            all_captions = caption_text_parts
        if transcript_segments:
            video_doc = builder.build(video_id=video_id, source=str(vp), transcript_segments=transcript_segments, ocr_results=ocr_results, video_path=str(vp), duration=duration)
            if all_captions:
                from src.video.video_document_builder import VideoChunk
                video_doc.chunks.append(VideoChunk(timestamp_start=0.0, timestamp_end=round(duration, 2), caption_text='\n'.join(all_captions), metadata={'chunk_type': 'blip_captions', 'caption_count': len(all_captions)}))
        else:
            video_doc = builder.build_from_captions(video_id=video_id, source=str(vp), captions=all_captions, video_path=str(vp), duration=duration, ocr_results=ocr_results)
        builder.save(video_doc)
    except Exception as exc:
        logger.error('VideoDocumentBuilder failed for %s: %s', vp.name, exc)
        return None
    text = video_doc.text
    if not text.strip():
        logger.warning('Video %s produced empty text after processing', vp.name)
        return None
    return RawDocument(doc_key=video_id, source=str(vp), text=text, metadata={'file_type': 'video', 'file_name': vp.name, 'source_directory': str(vp.parent), 'modality': 'video', 'duration': round(duration, 2), 'frame_count': len(frame_paths), 'ocr_frames_with_text': len(ocr_results), 'transcript_segments': len(transcript_segments), 'caption_count': len(frame_captions), 'caption_interval': caption_interval, 'video_doc_id': video_doc.doc_id, 'video_frame_paths': frame_paths})

def _handle_funsd_record(record: Dict, idx: int, split: str, image_dir: str) -> RawDocument:
    words = record.get('words', [])
    text = ' '.join(words) if words else ''
    image_path = None
    if record.get('image') is not None:
        fname = f'funsd_{split}_{idx:05d}.png'
        image_path = _extract_image(record['image'], os.path.join(image_dir, 'funsd', fname))
    return RawDocument(doc_key=f'funsd_{split}_{idx:05d}', source='funsd', text=text, image_path=image_path, metadata={'dataset': 'funsd', 'split': split, 'record_index': idx, 'word_count': len(words), 'has_bboxes': bool(record.get('bboxes'))})

def _handle_docvqa_record(record: Dict, idx: int, image_dir: str) -> RawDocument:
    image_path = None
    if record.get('image') is not None:
        fname = f'docvqa_{idx:05d}.png'
        image_path = _extract_image(record['image'], os.path.join(image_dir, 'docvqa', fname))
    return RawDocument(doc_key=f'docvqa_{idx:05d}', source='docvqa', text='', image_path=image_path, metadata={'dataset': 'docvqa', 'record_index': idx, 'question': record.get('question', ''), 'answers': record.get('answers', []), 'question_id': record.get('questionId', ''), 'doc_id': record.get('docId')})

def _handle_rvl_cdip_record(record: Dict, idx: int, image_dir: str) -> RawDocument:
    image_path = None
    if record.get('image') is not None:
        fname = f'rvl_cdip_{idx:05d}.png'
        image_path = _extract_image(record['image'], os.path.join(image_dir, 'rvl_cdip', fname))
    label_names = ['letter', 'form', 'email', 'handwritten', 'advertisement', 'scientific_report', 'scientific_publication', 'specification', 'file_folder', 'news_article', 'budget', 'invoice', 'presentation', 'questionnaire', 'resume', 'memo']
    label_int = record.get('label', -1)
    label_str = label_names[label_int] if 0 <= label_int < 16 else 'unknown'
    return RawDocument(doc_key=f'rvl_cdip_{idx:05d}', source='rvl_cdip', text='', image_path=image_path, metadata={'dataset': 'rvl_cdip', 'record_index': idx, 'label': label_int, 'label_name': label_str})
_RECORD_HANDLERS = {'funsd': _handle_funsd_record, 'docvqa': _handle_docvqa_record, 'rvl_cdip': _handle_rvl_cdip_record}

class DatasetLoader:

    def __init__(self, raw_data_root: str, image_cache_dir: str):
        self.raw_data_root = Path(raw_data_root)
        self.image_cache_dir = Path(image_cache_dir)
        if not self.raw_data_root.exists():
            raise FileNotFoundError(f'Raw data root does not exist: {self.raw_data_root}')
        self.image_cache_dir.mkdir(parents=True, exist_ok=True)

    def load_dataset(self, name: str, max_records: int=0, splits: Optional[List[str]]=None) -> List[RawDocument]:
        dataset_path = self.raw_data_root / name
        if not dataset_path.exists():
            raise FileNotFoundError(f'Dataset directory not found: {dataset_path}')
        logger.info("Loading dataset '%s' from %s", name, dataset_path)
        ds = _load_arrow_dataset(str(dataset_path))
        handler = _RECORD_HANDLERS.get(name)
        if handler is None:
            logger.warning("No handler registered for dataset '%s'.  Returning empty list. Register a handler in loader.py.", name)
            return []
        return self._iterate_dataset(ds, name, handler, max_records, splits)

    def load_directory(self, dir_path: str, max_records: int=0) -> List[RawDocument]:
        root = Path(dir_path)
        if not root.exists():
            raise FileNotFoundError(f'Directory not found: {root}')
        cache_mgr = CacheManager()
        documents: List[RawDocument] = []
        for file_path in sorted(root.rglob('*')):
            if not file_path.is_file():
                continue
            if max_records > 0 and len(documents) >= max_records:
                break
            ext = file_path.suffix.lower()
            rel = str(file_path.relative_to(root))
            cached_data = cache_mgr.get_cached(str(file_path))
            if cached_data is not None:
                meta = cached_data.get('metadata', {})
                meta['is_cached'] = True
                meta['cached_chunks'] = cached_data.get('chunks', [])
                meta['cached_clip_data'] = cached_data.get('clip_data', {})
                doc = RawDocument(doc_key=cached_data.get('doc_key', str(file_path)), source=cached_data.get('source', str(file_path)), text=cached_data.get('text', ''), image_path=cached_data.get('image_path'), metadata=meta)
                logger.info('Loaded cached document for: %s', file_path.name)
                documents.append(doc)
                continue
            if ext in _FILE_EXTRACTORS:
                try:
                    text = _FILE_EXTRACTORS[ext](str(file_path))
                except Exception as exc:
                    logger.warning('Failed to extract text from %s: %s', file_path, exc)
                    text = ''
                meta: Dict = {'file_type': ext, 'file_name': file_path.name, 'source_directory': str(root)}
                if ext == '.docx':
                    meta.update(_extract_docx_metadata(str(file_path)))
                elif ext == '.pptx':
                    meta.update(_extract_pptx_metadata(str(file_path)))
                documents.append(RawDocument(doc_key=rel, source=str(file_path), text=text, metadata=meta))
            elif ext in _IMAGE_EXTENSIONS:
                documents.append(RawDocument(doc_key=rel, source=str(file_path), text='', image_path=str(file_path.resolve()), metadata={'file_type': ext, 'file_name': file_path.name, 'source_directory': str(root)}))
            elif ext in _VIDEO_EXTENSIONS:
                try:
                    import yaml as _yaml
                    _settings_path = Path(__file__).resolve().parent.parent.parent / 'configs' / 'settings.yaml'
                    _settings = {}
                    if _settings_path.exists():
                        _settings = _yaml.safe_load(_settings_path.read_text()) or {}
                except Exception:
                    _settings = {}
                if 'video' not in _settings:
                    _settings['video'] = {}
                _settings['video']['enable_whisper'] = True
                logger.info('Processing video file: %s', file_path)
                video_doc = _process_video_file(str(file_path), settings=_settings)
                if video_doc is not None:
                    documents.append(video_doc)
                else:
                    logger.warning('Video processing returned no text for: %s', file_path)
            else:
                logger.debug('Skipping unsupported file: %s', file_path)
        logger.info("Directory '%s': %d documents loaded", dir_path, len(documents))
        return documents

    def load_path(self, path: str, max_records: int=0) -> List[RawDocument]:
        p = Path(path)
        if not p.exists():
            raise FileNotFoundError(f'Path not found: {p}')
        if p.is_dir():
            return self.load_directory(str(p), max_records=max_records)
        cache_mgr = CacheManager()
        cached_data = cache_mgr.get_cached(str(p))
        if cached_data is not None:
            meta = cached_data.get('metadata', {})
            meta['is_cached'] = True
            meta['cached_chunks'] = cached_data.get('chunks', [])
            meta['cached_clip_data'] = cached_data.get('clip_data', {})
            doc = RawDocument(doc_key=cached_data.get('doc_key', str(p)), source=cached_data.get('source', str(p)), text=cached_data.get('text', ''), image_path=cached_data.get('image_path'), metadata=meta)
            logger.info('Loaded cached document for: %s', p.name)
            return [doc]
        ext = p.suffix.lower()
        if ext in _FILE_EXTRACTORS:
            try:
                text = _FILE_EXTRACTORS[ext](str(p))
            except Exception as exc:
                logger.warning('Failed to extract text from %s: %s', p, exc)
                text = ''
            meta: Dict = {'file_type': ext, 'file_name': p.name, 'source_directory': str(p.parent)}
            if ext == '.docx':
                meta.update(_extract_docx_metadata(str(p)))
            elif ext == '.pptx':
                meta.update(_extract_pptx_metadata(str(p)))
            return [RawDocument(doc_key=p.name, source=str(p), text=text, metadata=meta)]
        if ext in _IMAGE_EXTENSIONS:
            return [RawDocument(doc_key=p.name, source=str(p), text='', image_path=str(p.resolve()), metadata={'file_type': ext, 'file_name': p.name, 'source_directory': str(p.parent)})]
        if ext in _VIDEO_EXTENSIONS:
            try:
                import yaml as _yaml
                _settings_path = Path(__file__).resolve().parent.parent.parent / 'configs' / 'settings.yaml'
                _settings = {}
                if _settings_path.exists():
                    _settings = _yaml.safe_load(_settings_path.read_text()) or {}
            except Exception:
                _settings = {}
            if 'video' not in _settings:
                _settings['video'] = {}
            _settings['video']['enable_whisper'] = True
            logger.info('Processing video file: %s', p)
            video_doc = _process_video_file(str(p), settings=_settings)
            if video_doc is not None:
                return [video_doc]
            logger.warning('Video processing returned no text for: %s', p)
            return []
        logger.warning('Unsupported file type: %s', p)
        return []

    def list_available_datasets(self) -> List[str]:
        if not self.raw_data_root.exists():
            return []
        return sorted((d.name for d in self.raw_data_root.iterdir() if d.is_dir()))

    def _iterate_dataset(self, ds, name: str, handler, max_records: int, splits: Optional[List[str]]=None) -> List[RawDocument]:
        from datasets import DatasetDict
        documents: List[RawDocument] = []
        image_dir = str(self.image_cache_dir)
        if isinstance(ds, DatasetDict):
            available_splits = list(ds.keys())
            if splits is not None:
                selected = []
                for s in splits:
                    if s in ds:
                        selected.append(s)
                    else:
                        logger.warning("Requested split '%s' not found in dataset '%s'. Available splits: %s", s, name, available_splits)
                target_splits = selected
            else:
                target_splits = available_splits
            for split_name in target_splits:
                split_ds = ds[split_name]
                count = 0
                for idx, record in enumerate(split_ds):
                    if max_records > 0 and len(documents) >= max_records:
                        break
                    doc = handler(record, idx, split_name, image_dir)
                    documents.append(doc)
                    count += 1
                logger.info("  split '%s': loaded %d records", split_name, count)
        else:
            for idx, record in enumerate(ds):
                if max_records > 0 and len(documents) >= max_records:
                    break
                doc = handler(record, idx, image_dir)
                documents.append(doc)
        logger.info("Dataset '%s': %d total documents loaded", name, len(documents))
        return documents
